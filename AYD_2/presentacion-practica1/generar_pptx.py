# -*- coding: utf-8 -*-
"""
Genera la presentacion de la Practica 1 (AutoRent Express S.A.) en PowerPoint.
Curso: Analisis y Diseno de Sistemas 2 - USAC. Grupo 2, seccion A.

Uso:  python generar_pptx.py
Salida: Presentacion_Practica1_AutoRent_G2.pptx (misma carpeta)
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------
# Paleta y tipografia
# --------------------------------------------------------------------------
CYAN = RGBColor(0x29, 0xAB, 0xE2)   # azul del enunciado
CYAN_D = RGBColor(0x0E, 0x7A, 0xAD)
NAVY = RGBColor(0x10, 0x22, 0x38)
NAVY_2 = RGBColor(0x1B, 0x35, 0x52)
GOLD = RGBColor(0xE0, 0x9F, 0x1E)
GREEN = RGBColor(0x2E, 0x9E, 0x6B)
RED = RGBColor(0xC0, 0x39, 0x3D)
VIOLET = RGBColor(0x6C, 0x4E, 0xB3)
INK = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF2, 0xF7, 0xFB)
LINE = RGBColor(0xD8, 0xE3, 0xED)

FONT = "Segoe UI"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]

_page = {"n": 0}


# --------------------------------------------------------------------------
# Helpers
# --------------------------------------------------------------------------
def rect(slide, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.text_frame.text = ""
    return sh


def txt(slide, x, y, w, h, text, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT,
        font=FONT, italic=False, anchor=MSO_ANCHOR.TOP, spacing=1.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.alignment = align
        p.line_spacing = spacing
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = color
            r.font.name = font
    return box


def bullets(slide, x, y, w, h, items, size=15, color=INK, bullet_color=CYAN, gap=8,
            bold_lead=True):
    """items: lista de str o (lead, resto)."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.15
        r0 = p.add_run()
        r0.text = "▪  "
        r0.font.size = Pt(size)
        r0.font.color.rgb = bullet_color
        r0.font.bold = True
        r0.font.name = FONT
        if isinstance(it, tuple):
            lead, rest = it
            r1 = p.add_run()
            r1.text = lead
            r1.font.size = Pt(size)
            r1.font.bold = bold_lead
            r1.font.color.rgb = color
            r1.font.name = FONT
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
            r2.font.name = FONT
        else:
            r1 = p.add_run()
            r1.text = it
            r1.font.size = Pt(size)
            r1.font.color.rgb = color
            r1.font.name = FONT
    return box


def table(slide, x, y, w, data, col_w, row_h=0.34, head_h=0.4, size=11.5,
          head_fill=NAVY, head_color=WHITE, zebra=LIGHT, first_bold=False,
          head_size=None):
    rows, cols = len(data), len(data[0])
    gf = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w),
                                Inches(head_h + row_h * (rows - 1)))
    tbl = gf.table
    tbl.first_row = True
    tbl.horz_banding = False
    total = sum(col_w)
    for j, cw in enumerate(col_w):
        tbl.columns[j].width = Emu(int(Inches(w) * cw / total))
    tbl.rows[0].height = Inches(head_h)
    for i in range(1, rows):
        tbl.rows[i].height = Inches(row_h)

    for i, row in enumerate(data):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = str(val)
            c.margin_left = Inches(0.09)
            c.margin_right = Inches(0.07)
            c.margin_top = Inches(0.03)
            c.margin_bottom = Inches(0.03)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            if i == 0:
                c.fill.fore_color.rgb = head_fill
            else:
                c.fill.fore_color.rgb = WHITE if (i % 2 == 1) else zebra
            for p in c.text_frame.paragraphs:
                p.line_spacing = 1.0
                for r in p.runs:
                    r.font.name = FONT
                    r.font.size = Pt(head_size or size) if i == 0 else Pt(size)
                    r.font.bold = (i == 0) or (first_bold and j == 0)
                    r.font.color.rgb = head_color if i == 0 else INK
    return gf


def slide_base(title, kicker=None, rule=True):
    """Lamina de contenido con barra superior, kicker y titulo."""
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.333, 7.5, WHITE)
    rect(s, 0, 0, 13.333, 0.13, CYAN)
    if kicker:
        txt(s, 0.72, 0.42, 10.0, 0.28, kicker.upper(), 11.5, True, CYAN)
        txt(s, 0.72, 0.72, 11.9, 0.55, title, 27, True, NAVY)
        ry = 1.36
    else:
        txt(s, 0.72, 0.55, 11.9, 0.6, title, 28, True, NAVY)
        ry = 1.24
    if rule:
        rect(s, 0.72, ry, 1.05, 0.05, CYAN)
    _page["n"] += 1
    txt(s, 11.9, 6.92, 0.75, 0.3, f"{_page['n']:02d}", 11, True, MUTED, PP_ALIGN.RIGHT)
    txt(s, 0.72, 6.92, 8.0, 0.3, "AutoRent Express S.A.  ·  Practica 1  ·  Grupo 2",
        10, False, MUTED)
    return s


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def card(slide, x, y, w, h, title, body, accent=CYAN, tsize=15, bsize=12.5, icon=None):
    rect(slide, x, y, w, h, LIGHT, LINE, 0.75)
    rect(slide, x, y, 0.055, h, accent)
    ty = y + 0.18
    if icon:
        txt(slide, x + 0.22, ty, w - 0.4, 0.35, icon, 18, True, accent)
        ty += 0.42
    txt(slide, x + 0.22, ty, w - 0.42, 0.32, title, tsize, True, NAVY)
    txt(slide, x + 0.22, ty + 0.36, w - 0.42, h - (ty - y) - 0.5, body, bsize, False, INK,
        spacing=1.15)


def code_box(slide, x, y, w, h, code, size=11.5, fill=NAVY, color=RGBColor(0xE6, 0xF4, 0xFC)):
    rect(slide, x, y, w, h, fill, None)
    box = slide.shapes.add_textbox(Inches(x + 0.16), Inches(y + 0.12),
                                   Inches(w - 0.32), Inches(h - 0.24))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, ln in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.line_spacing = 1.2
        for r in p.runs:
            r.font.name = MONO
            r.font.size = Pt(size)
            r.font.color.rgb = color


def chip(slide, x, y, w, h, text, fill, tcolor=WHITE, size=11.5):
    rect(slide, x, y, w, h, fill, None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(slide, x, y + (h - 0.22) / 2, w, 0.24, text, size, True, tcolor, PP_ALIGN.CENTER)


def arrow(slide, x, y, w, h=0.3, color=CYAN):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.shadow.inherit = False
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.text_frame.text = ""


# ==========================================================================
# 01 - Portada
# ==========================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 0, 0, 13.333, 0.22, CYAN)
rect(s, 0, 7.28, 13.333, 0.22, CYAN)
rect(s, 7.9, 0.22, 5.44, 7.06, NAVY_2)
txt(s, 1.0, 1.55, 7.0, 0.32, "UNIVERSIDAD DE SAN CARLOS DE GUATEMALA", 12.5, True, CYAN)
txt(s, 1.0, 1.9, 7.0, 0.3, "Facultad de Ingenieria  ·  Escuela de Ciencias y Sistemas",
    12, False, RGBColor(0xB6, 0xC6, 0xD6))
rect(s, 1.0, 2.5, 1.3, 0.07, CYAN)
txt(s, 1.0, 2.85, 7.0, 1.1, "AutoRent Express S.A.", 46, True, WHITE)
txt(s, 1.0, 3.95, 7.0, 0.75,
    "Plataforma digital para la gestion y publicacion\nde vehiculos en alquiler",
    17, False, RGBColor(0xCB, 0xDA, 0xE8), spacing=1.2)
rect(s, 1.0, 5.0, 6.2, 0.02, RGBColor(0x2E, 0x4A, 0x68))
txt(s, 1.0, 5.25, 7.0, 0.3, "Analisis y Diseno de Sistemas 2  ·  Practica 1  ·  Seccion A",
    14, True, WHITE)
txt(s, 1.0, 5.62, 7.0, 0.3, "Grupo 2  ·  Segundo semestre 2026", 13, False,
    RGBColor(0x9F, 0xB4, 0xC8))

txt(s, 8.5, 1.55, 4.3, 0.3, "INTEGRANTES", 12, True, CYAN)
rect(s, 8.5, 1.9, 0.8, 0.05, GOLD)
integrantes = [
    ("201801521", "Jemima Solmaira Chavajay"),
    ("201900131", "Elias Abraham Vasquez Soto"),
    ("201905884", "Santiago Julian Barrera Reyes"),
    ("202112395", "Evelyn Marisol Pumay Soy"),
    ("201701078", "Byron Josue Par Rancho"),
]
yy = 2.25
for carne, nombre in integrantes:
    txt(s, 8.5, yy, 1.35, 0.26, carne, 12, True, GOLD, font=MONO)
    txt(s, 9.85, yy, 3.1, 0.26, nombre, 12.5, False, WHITE)
    yy += 0.52
rect(s, 8.5, 5.15, 4.3, 0.02, RGBColor(0x2E, 0x4A, 0x68))
txt(s, 8.5, 5.4, 4.3, 0.28, "Repositorio", 11, True, CYAN)
txt(s, 8.5, 5.68, 4.5, 0.3, "AYD2_A_2S2026_PRACTICA1_G2", 12, True, WHITE, font=MONO)
txt(s, 8.5, 6.05, 4.3, 0.28, "Django 6.0.7  ·  PostgreSQL  ·  Git Flow", 11.5,
    False, RGBColor(0x9F, 0xB4, 0xC8))
notes(s, "Presentacion de la Practica 1. Presentarse, decir grupo y seccion, y anunciar "
         "que se recorre el enunciado completo: framework, patron, demo, git-flow y "
         "documentacion.")

# ==========================================================================
# 02 - Agenda
# ==========================================================================
s = slide_base("Agenda", "Recorrido de la practica")
items = [
    ("01", "Contexto y enunciado", "Que pidio AutoRent Express y quienes son los actores"),
    ("02", "Framework elegido", "Django + PostgreSQL y su justificacion"),
    ("03", "Patrones de diseno", "State, Strategy y Factory Method con UML"),
    ("04", "Implementacion del demo", "Arquitectura MTV, apps y modelo de datos"),
    ("05", "Requerimientos", "14 funcionales, 15 casos de uso y 10 no funcionales"),
    ("06", "Git Flow", "Ramas, features y formato de commits"),
    ("07", "Entregables y demo", "Documentacion en markdown y recorrido en vivo"),
]
yy = 1.72
for num, t, d in items:
    rect(s, 0.72, yy, 0.52, 0.52, CYAN)
    txt(s, 0.72, yy + 0.13, 0.52, 0.3, num, 15, True, WHITE, PP_ALIGN.CENTER, font=MONO)
    txt(s, 1.45, yy + 0.03, 4.0, 0.3, t, 16, True, NAVY)
    txt(s, 1.45, yy + 0.29, 8.5, 0.3, d, 13, False, MUTED)
    yy += 0.72
rect(s, 10.35, 1.72, 2.25, 4.35, LIGHT, LINE, 0.75)
txt(s, 10.6, 1.95, 1.8, 0.3, "EN UNA FRASE", 10.5, True, CYAN)
txt(s, 10.6, 2.3, 1.8, 3.4,
    "Un framework, tres patrones y un ciclo de renta que ningun actor puede avanzar solo.",
    14, True, NAVY, spacing=1.25)
notes(s, "No leer la agenda entera. Marcar los tres bloques gruesos: framework y patrones, "
         "implementacion, y proceso de trabajo (git-flow y documentacion).")

# ==========================================================================
# 03 - Competencias y objetivos
# ==========================================================================
s = slide_base("Competencias y objetivos", "Lo que pide el enunciado")
txt(s, 0.72, 1.62, 5.6, 0.3, "COMPETENCIA", 11.5, True, CYAN)
rect(s, 0.72, 1.95, 5.6, 0.95, LIGHT, LINE, 0.75)
txt(s, 0.95, 2.12, 5.15, 0.72,
    "Identificar y describir conceptos de frameworks, aplicando casos de uso y "
    "caracteristicas comunes para seleccionar el framework adecuado al proyecto.",
    12.5, False, INK, spacing=1.2)
txt(s, 0.72, 3.08, 5.6, 0.3, "OBJETIVOS GENERALES", 11.5, True, CYAN)
bullets(s, 0.72, 3.42, 5.6, 1.8, [
    "Familiarizarse con los frameworks de desarrollo y los patrones de diseno",
    "Comprender la importancia y la organizacion del trabajo en equipo",
    "Establecer un flujo de trabajo estructurado que facilite la colaboracion",
], 13, gap=6)
txt(s, 0.72, 5.4, 5.6, 0.3, "DATOS DE LA PRACTICA", 11.5, True, CYAN)
datos = [("Ponderacion", "5 puntos"), ("Carga estimada", "30 horas"),
         ("Entrega", "5 de agosto 2026, por UEDI"), ("Calificacion", "8 de agosto 2026")]
for i, (k, v) in enumerate(datos):
    dx = 0.72 + (i % 2) * 2.86
    dy = 5.74 + (i // 2) * 0.58
    rect(s, dx, dy, 2.74, 0.52, LIGHT, LINE, 0.6)
    txt(s, dx + 0.15, dy + 0.05, 2.4, 0.22, k.upper(), 9.5, True, CYAN)
    txt(s, dx + 0.15, dy + 0.26, 2.5, 0.24, v, 11.5, False, INK)
txt(s, 7.05, 1.62, 5.55, 0.3, "OBJETIVOS ESPECIFICOS", 11.5, True, GOLD)
esp = [
    ("Utilizar un framework e implementar patrones de diseno",
     "Django 6.0.7 + State, Strategy y Factory Method"),
    ("Implementar git-flow para el control de versiones",
     "main, develop y 15 ramas feature con commits #carne: mensaje"),
    ("Identificar requerimientos funcionales y no funcionales",
     "14 RF con plantilla IEEE 830 y 10 RNF medibles"),
]
yy = 2.0
for t, d in esp:
    rect(s, 7.05, yy, 5.55, 1.12, WHITE, LINE, 0.75)
    rect(s, 7.05, yy, 0.055, 1.12, GOLD)
    txt(s, 7.3, yy + 0.18, 5.1, 0.3, t, 13.5, True, NAVY)
    txt(s, 7.3, yy + 0.55, 5.1, 0.45, "Como lo cumplimos:  " + d, 12, False, MUTED,
        spacing=1.15)
    yy += 1.27
rect(s, 7.05, 5.82, 5.55, 0.55, NAVY)
txt(s, 7.3, 5.95, 5.1, 0.3, "Los tres objetivos especificos se cubren en esta presentacion.",
    12.5, True, WHITE)
notes(s, "Enlazar cada objetivo especifico con la evidencia concreta del proyecto. Esto "
         "adelanta el mapa de la defensa.")

# ==========================================================================
# 04 - Enunciado
# ==========================================================================
s = slide_base("El problema de negocio", "Descripcion / Enunciado")
rect(s, 0.72, 1.62, 11.9, 1.0, LIGHT, LINE, 0.75)
rect(s, 0.72, 1.62, 0.055, 1.0, CYAN)
txt(s, 1.0, 1.82, 11.4, 0.7,
    "AutoRent Express S.A., dedicada a la renta de vehiculos, necesita una plataforma "
    "digital para administrar usuarios, vehiculos y planes de renta, con roles y permisos "
    "diferenciados y con interfaces administrativas para gestionar la flota.",
    14.5, False, INK, spacing=1.2)
txt(s, 0.72, 2.92, 6.0, 0.3, "QUE DEBE PERMITIR LA PLATAFORMA", 11.5, True, CYAN)
cols = [
    ("Registro y autenticacion", "Clientes, agentes y administradores entran con correo "
     "electronico y contrasena.", CYAN),
    ("Catalogo y favoritos", "El cliente busca vehiculos, guarda favoritos y consulta la "
     "ficha completa.", VIOLET),
    ("Ciclo de la reserva", "Solicitud, propuesta, aceptacion o rechazo, entrega y "
     "finalizacion.", GOLD),
    ("Administracion", "Alta y baja de agentes, gestion de vehiculos y tipos, dashboard "
     "operativo.", GREEN),
]
x = 0.72
for t, d, c in cols:
    card(s, x, 3.3, 2.87, 1.65, t, d, c, 14, 12)
    x += 3.03
txt(s, 0.72, 5.25, 6.0, 0.3, "LO DIFICIL NO ES LISTAR CARROS", 11.5, True, GOLD)
rect(s, 0.72, 5.6, 11.9, 0.95, NAVY)
txt(s, 1.0, 5.78, 11.4, 0.65,
    "El nucleo del problema es la negociacion entre cliente y agente para cerrar una renta, "
    "y el precio, que depende del plan contratado. De ahi salieron los patrones: nadie puede "
    "avanzar la reserva solo, y el total no se calcula igual en todos los planes.",
    13.5, False, WHITE, spacing=1.2)
notes(s, "Decir el enunciado con las propias palabras. El gancho es la ultima caja: el "
         "problema real es el ida y vuelta cliente-agente, no el CRUD.")

# ==========================================================================
# 05 - Actores
# ==========================================================================
s = slide_base("Actores y alcance por rol", "Quien hace que")
data = [
    ["Actor", "Responsabilidad", "Alcance implementado"],
    ["Cliente", "Busca, guarda favoritos y solicita reservas",
     "Registro con nombre, correo unico, contrasena, licencia y fotografia; "
     "actualizacion de perfil; solicitud con fecha, hora, direccion y dias; "
     "aceptar o rechazar la propuesta"],
    ["Agente de renta", "Publica y mantiene la flota",
     "CRUD de vehiculos (marca, modelo, ano, placa, precio/dia, descripcion, fotos, "
     "pasajeros, transmision, combustible); tipos de vehiculo; dashboard; propuestas "
     "de horario"],
    ["Administrador", "Gobierna a los agentes",
     "Crear y dar de baja agentes de renta"],
    ["Sistema", "Valida, persiste y calcula",
     "Aplica transiciones de estado, calcula tarifas y hace cumplir los constraints "
     "del esquema"],
]
table(s, 0.72, 1.62, 11.9, data, [1.5, 2.6, 7.5], row_h=0.92, head_h=0.42, size=12.5,
      first_bold=True)
rect(s, 0.72, 6.0, 11.9, 0.62, LIGHT, LINE, 0.75)
txt(s, 0.95, 6.16, 11.4, 0.3,
    "Autenticacion unica para los tres roles: correo electronico y contrasena. "
    "El rol decide a que pantallas se llega.", 13, True, NAVY)
notes(s, "Aclarar que el actor Sistema no esta en el enunciado: lo agregamos porque en los "
         "casos de uso expandidos hay pasos que ejecuta la aplicacion sola.")

# ==========================================================================
# 06 - Flujo de la reserva
# ==========================================================================
s = slide_base("El recorrido de una renta", "Flujo principal del negocio")
pasos = [
    ("1", "Solicita", "Cliente", "Vehiculo, fecha,\nhora, direccion y dias", CYAN),
    ("2", "Propone", "Agente", "Horario de entrega\ny plan de renta", VIOLET),
    ("3", "Acepta o\nrechaza", "Cliente", "Si rechaza, el sistema\nexige un motivo", GOLD),
    ("4", "Entrega", "Agente", "La reserva pasa\na EN RENTA", GREEN),
    ("5", "Finaliza", "Agente", "Devolucion y cierre\nde la renta", NAVY),
]
x = 0.72
for num, t, actor, d, c in pasos:
    rect(s, x, 1.72, 2.05, 2.5, WHITE, LINE, 0.75)
    rect(s, x, 1.72, 2.05, 0.06, c)
    rect(s, x + 0.78, 1.95, 0.5, 0.5, c, shape=MSO_SHAPE.OVAL)
    txt(s, x + 0.78, 2.07, 0.5, 0.3, num, 15, True, WHITE, PP_ALIGN.CENTER, font=MONO)
    txt(s, x + 0.15, 2.58, 1.75, 0.55, t, 15, True, NAVY, PP_ALIGN.CENTER, spacing=1.05)
    txt(s, x + 0.15, 3.16, 1.75, 0.26, actor.upper(), 10.5, True, c, PP_ALIGN.CENTER)
    txt(s, x + 0.12, 3.5, 1.81, 0.6, d, 11, False, MUTED, PP_ALIGN.CENTER, spacing=1.15)
    if x < 10.0:
        arrow(s, x + 2.12, 2.85, 0.32, 0.26, LINE)
    x += 2.48
rect(s, 0.72, 4.55, 11.9, 0.75, LIGHT, LINE, 0.75)
rect(s, 0.72, 4.55, 0.055, 0.75, CYAN)
txt(s, 1.0, 4.73, 11.4, 0.4,
    "Ninguno de los dos actores puede avanzar solo: cliente y agente se van pasando el turno. "
    "Esa alternancia es exactamente lo que modela el patron State.", 13.5, False, INK)
rect(s, 0.72, 5.55, 11.9, 1.0, NAVY)
txt(s, 1.0, 5.72, 3.3, 0.3, "EL PASO 3 ES EL CRITICO", 11.5, True, CYAN)
txt(s, 1.0, 6.02, 11.4, 0.4,
    "Cuando el cliente acepta, se disparan los tres patrones en cadena: State valida la "
    "transicion, Factory entrega el algoritmo y Strategy calcula el precio. Es el unico "
    "momento en que se calcula dinero.", 13, False, WHITE)
notes(s, "Contar el flujo como historia, sin leer las tarjetas. Rematar con la franja "
         "inferior: los tres patrones se concentran en el paso 3.")

# ==========================================================================
# 07 - Framework
# ==========================================================================
s = slide_base("Framework seleccionado", "Instruccion 1a")
rect(s, 0.72, 1.62, 5.5, 2.05, NAVY)
txt(s, 1.0, 1.9, 5.0, 0.4, "Django 6.0.7", 34, True, WHITE)
txt(s, 1.0, 2.45, 5.0, 0.3, "Framework web full-stack en Python", 14, False,
    RGBColor(0xB6, 0xC6, 0xD6))
rect(s, 1.0, 2.9, 4.9, 0.02, RGBColor(0x2E, 0x4A, 0x68))
txt(s, 1.0, 3.05, 5.0, 0.3, "Persistencia:  PostgreSQL", 13.5, True, CYAN)
tecno = [
    ("MTV", "Model - Template - View", CYAN),
    ("ORM", "Modelos mapeados al DDL", VIOLET),
    ("Auth", "Sesiones y hash de contrasena", GREEN),
    ("CSRF", "Middleware de seguridad", GOLD),
]
x, yy = 6.55, 1.62
for t, d, c in tecno:
    rect(s, x, yy, 2.95, 0.95, LIGHT, LINE, 0.75)
    rect(s, x, yy, 0.055, 0.95, c)
    txt(s, x + 0.2, yy + 0.16, 2.6, 0.28, t, 14, True, NAVY)
    txt(s, x + 0.2, yy + 0.5, 2.65, 0.3, d, 11.5, False, MUTED)
    x += 3.1
    if x > 12.0:
        x, yy = 6.55, yy + 1.1
txt(s, 0.72, 3.95, 6.0, 0.3, "POR QUE DJANGO Y NO OTRO STACK", 11.5, True, CYAN)
data = [
    ["Criterio", "Argumento"],
    ["Productividad", "El patron MTV permitio construir frontend, backend y persistencia "
     "dentro de un solo framework, en el tiempo del demo"],
    ["Autenticacion", "Sesiones, hash de contrasenas y AUTH_USER_MODEL personalizado con "
     "login por correo, sin escribir el modulo desde cero"],
    ["ORM sobre DDL previo", "El esquema se diseno primero en SQL; los modelos lo reflejan "
     "con db_table y managed = False, respetando las 11 tablas"],
    ["Seguridad", "Middleware CSRF, validadores de contrasena y secretos por variables de "
     "entorno"],
    ["Expresion de patrones", "Permite ubicar State, Strategy y Factory en clases Python "
     "del dominio, sin acoplarlos a la capa de presentacion"],
]
table(s, 0.72, 4.3, 11.9, data, [2.5, 9.4], row_h=0.44, head_h=0.38, size=12,
      first_bold=True)
notes(s, "Cerrar diciendo que se descarto un stack de solo frontend con API ad hoc: por el "
         "tiempo del demo y por la necesidad de auth, ORM y admin ya resueltos.")

# ==========================================================================
# 08 - Patrones: panoramica
# ==========================================================================
s = slide_base("Tres patrones, no uno", "Instruccion 1b")
txt(s, 0.72, 1.6, 11.9, 0.35,
    "El enunciado pedia elegir un patron. El dominio pidio tres: el ciclo de vida de la "
    "reserva y el calculo del precio son problemas distintos.", 14, False, MUTED)
pats = [
    ("State", "Comportamiento", "Transiciones legales del ciclo de vida de la reserva",
     "reservas/estados.py", "Elimino las cadenas if/elif sobre el campo estado", CYAN),
    ("Strategy", "Comportamiento", "Algoritmos intercambiables de calculo de tarifa",
     "reservas/calculos.py", "Una clase por plan: diaria, semanal y mensual", VIOLET),
    ("Factory Method", "Creacional", "Instanciar la estrategia correcta segun el plan "
     "persistido", "PlanRenta.calculo()", "Conecta el discriminador de la base con el "
     "objeto de comportamiento", GOLD),
]
x = 0.72
for nombre, fam, prob, loc, ben, c in pats:
    rect(s, x, 2.15, 3.83, 3.55, WHITE, LINE, 0.9)
    rect(s, x, 2.15, 3.83, 0.07, c)
    txt(s, x + 0.25, 2.42, 3.3, 0.26, fam.upper(), 10.5, True, c)
    txt(s, x + 0.25, 2.72, 3.4, 0.42, nombre, 24, True, NAVY)
    rect(s, x + 0.25, 3.28, 0.7, 0.04, c)
    txt(s, x + 0.25, 3.48, 3.35, 0.75, prob, 13, False, INK, spacing=1.2)
    rect(s, x + 0.25, 4.3, 3.35, 0.34, LIGHT)
    txt(s, x + 0.38, 4.38, 3.1, 0.3, loc, 11, True, CYAN_D, font=MONO)
    txt(s, x + 0.25, 4.82, 3.35, 0.75, ben, 12, False, MUTED, spacing=1.2)
    x += 4.03
rect(s, 0.72, 5.95, 11.9, 0.7, NAVY)
txt(s, 1.0, 6.13, 11.4, 0.35,
    "No son tres patrones sueltos: se activan en cadena en una sola operacion, cuando el "
    "cliente acepta la propuesta.", 13.5, True, WHITE)
notes(s, "Si preguntan por que tres y no uno: el dominio los pidio. Con solo State se "
         "cumplia, pero el precio habria quedado con condicionales dentro de Reserva.")

# ==========================================================================
# 09 - State
# ==========================================================================
s = slide_base("State — ciclo de vida de la reserva", "Patron de comportamiento")
txt(s, 0.72, 1.6, 5.6, 0.3, "EL PROBLEMA SIN EL PATRON", 11.5, True, RED)
rect(s, 0.72, 1.95, 5.6, 1.35, LIGHT, LINE, 0.75)
rect(s, 0.72, 1.95, 0.055, 1.35, RED)
txt(s, 0.98, 2.13, 5.2, 1.0,
    "Cinco operaciones (proponer, aceptar, rechazar, entregar, finalizar) y cada una con su "
    "propio if sobre el campo estado. La regla queda dispersa, agregar un estado obliga a "
    "tocar las cinco, y es facil olvidar una.", 12.5, False, INK, spacing=1.15)
txt(s, 0.72, 3.5, 5.6, 0.3, "LA SOLUCION", 11.5, True, GREEN)
bullets(s, 0.72, 3.85, 5.6, 2.1, [
    ("Reserva es el Context. ", "Persiste el codigo del estado y delega."),
    ("EstadoReserva es la clase abstracta. ", "Sus cinco metodos lanzan ValidationError "
     "por omision."),
    ("Cada estado concreto ", "sobreescribe unicamente las transiciones que le son legales."),
    ("Un estado terminal ", "son tres lineas: hereda cinco metodos que ya rechazan."),
], 12.5)
txt(s, 6.75, 1.6, 5.85, 0.3, "MAQUINA DE ESTADOS IMPLEMENTADA", 11.5, True, CYAN)
estados = [
    ("SOLICITADA", 1.98, CYAN), ("PROPUESTA", 2.62, VIOLET),
    ("ACEPTADA", 3.26, GREEN), ("EN_RENTA", 3.9, GOLD),
    ("FINALIZADA", 4.54, NAVY),
]
for nombre, yy, c in estados:
    chip(s, 6.75, yy, 2.35, 0.44, nombre, c, WHITE, 12)
trans = [
    ("agente propone", 2.05), ("cliente acepta", 2.69),
    ("agente entrega", 3.33), ("agente finaliza", 3.97),
]
for lbl, yy in trans:
    txt(s, 9.3, yy, 1.6, 0.28, "↓  " + lbl, 11.5, False, MUTED)
chip(s, 11.05, 2.62, 1.55, 0.44, "RECHAZADA", RED, WHITE, 11.5)
arrow(s, 10.7, 2.72, 0.3, 0.24, RED)
txt(s, 10.35, 2.26, 2.25, 0.26, "cliente rechaza", 10.5, False, RED, PP_ALIGN.CENTER)
rect(s, 6.75, 5.15, 5.85, 0.85, LIGHT, LINE, 0.75)
txt(s, 6.98, 5.3, 5.4, 0.6,
    "RECHAZADA y FINALIZADA son terminales. Cada transicion sincroniza ademas el estado del "
    "vehiculo, dentro de la misma clase de estado.", 12, False, INK, spacing=1.15)
rect(s, 0.72, 6.15, 11.9, 0.55, NAVY)
txt(s, 1.0, 6.28, 11.4, 0.3,
    "Resultado medible: cero condicionales sobre el estado dentro de la clase Reserva. "
    "Las reglas estan espejadas en los CHECK del DDL.", 13, True, WHITE)
notes(s, "El argumento fuerte: la clase base no declara metodos vacios, declara metodos que "
         "lanzan excepcion. El comportamiento por omision es 'esta transicion es ilegal'. "
         "Mencionar chk_motivo_si_rechazada y chk_horario_si_propuesta.")

# ==========================================================================
# 10 - State UML
# ==========================================================================
s = slide_base("State — diagrama de clases", "UML del patron")
rect(s, 0.72, 1.68, 3.4, 2.12, NAVY)
txt(s, 0.95, 1.84, 3.0, 0.26, "«Context»", 11, True, CYAN, PP_ALIGN.CENTER)
txt(s, 0.95, 2.1, 3.0, 0.32, "Reserva", 18, True, WHITE, PP_ALIGN.CENTER)
rect(s, 0.95, 2.51, 2.95, 0.02, RGBColor(0x2E, 0x4A, 0x68))
code = ("- estado: str\n+ proponer()\n+ aceptar()\n+ rechazar(motivo)\n"
        "+ entregar()\n+ finalizar()")
txt(s, 0.98, 2.6, 3.0, 1.15, code, 10, False, RGBColor(0xD4, 0xE6, 0xF4), font=MONO,
    spacing=1.1)
arrow(s, 4.25, 2.55, 0.75, 0.28, CYAN)
txt(s, 4.1, 2.2, 1.1, 0.26, "delega", 10.5, False, MUTED, PP_ALIGN.CENTER)
rect(s, 5.15, 1.68, 3.4, 2.12, CYAN_D)
txt(s, 5.38, 1.84, 3.0, 0.26, "«State» (abstracta)", 11, True,
    RGBColor(0xCB, 0xEC, 0xFA), PP_ALIGN.CENTER)
txt(s, 5.38, 2.1, 3.0, 0.32, "EstadoReserva", 17, True, WHITE, PP_ALIGN.CENTER)
rect(s, 5.38, 2.51, 2.95, 0.02, RGBColor(0x6B, 0xC2, 0xE4))
txt(s, 5.41, 2.6, 3.0, 1.15,
    "+ proponer()   → raise\n+ aceptar()    → raise\n+ rechazar()   → raise\n"
    "+ entregar()   → raise\n+ finalizar()  → raise", 10, False, WHITE, font=MONO,
    spacing=1.1)
rect(s, 8.9, 1.68, 3.72, 2.12, LIGHT, GOLD, 1.0)
txt(s, 9.12, 1.88, 3.3, 0.26, "DECISION DE DISENO", 10.5, True, GOLD)
txt(s, 9.12, 2.2, 3.3, 1.4,
    "La implementacion por defecto de las cinco operaciones rechaza. Cada estado concreto "
    "sobreescribe solamente las transiciones que le son legales; no escribe ni un if.",
    12, False, INK, spacing=1.2)
txt(s, 0.72, 3.95, 6.0, 0.3, "«ConcreteState»  ·  una clase por estado",
    11.5, True, CYAN)
concretos = [
    ("Solicitada", "proponer()", CYAN),
    ("Propuesta", "aceptar()\nrechazar()", VIOLET),
    ("Aceptada", "entregar()", GREEN),
    ("EnRenta", "finalizar()", GOLD),
    ("Rechazada", "(terminal)", RED),
    ("Finalizada", "(terminal)", NAVY),
]
x = 0.72
for nombre, m, c in concretos:
    rect(s, x, 4.32, 1.87, 1.25, WHITE, LINE, 0.75)
    rect(s, x, 4.32, 1.87, 0.05, c)
    txt(s, x + 0.12, 4.5, 1.65, 0.28, nombre, 12.5, True, NAVY, PP_ALIGN.CENTER)
    txt(s, x + 0.08, 4.85, 1.72, 0.6, m, 10.5, False, MUTED, PP_ALIGN.CENTER, font=MONO,
        spacing=1.2)
    x += 2.0
rect(s, 0.72, 5.85, 11.9, 0.85, LIGHT, LINE, 0.75)
rect(s, 0.72, 5.85, 0.055, 0.85, GREEN)
txt(s, 1.0, 6.0, 11.4, 0.6,
    "Aplicabilidad segun el GoF: usar State cuando varias operaciones contienen la misma "
    "estructura condicional sobre el estado del objeto. Nosotros teniamos cinco.",
    13, False, INK, spacing=1.15)
notes(s, "Diagrama tomado de docs/diagramas/patrones_uml.md. Senalar la nota ambar: es la "
         "decision de diseno que hace que un estado terminal sean tres lineas.")

# ==========================================================================
# 11 - Strategy
# ==========================================================================
s = slide_base("Strategy — calculo de tarifas", "Patron de comportamiento")
txt(s, 0.72, 1.6, 5.7, 0.3, "EL PROBLEMA", 11.5, True, RED)
rect(s, 0.72, 1.95, 5.7, 1.0, LIGHT, LINE, 0.75)
rect(s, 0.72, 1.95, 0.055, 1.0, RED)
txt(s, 0.98, 2.12, 5.3, 0.75,
    "Tres planes cobran distinto. Sin el patron, calcular_total tendria un if por plan, y "
    "agregar un plan quincenal obligaria a modificar una clase que no tiene nada que ver "
    "con formulas de precio.", 12.5, False, INK, spacing=1.15)
txt(s, 0.72, 3.15, 5.7, 0.3, "LAS TRES ESTRATEGIAS", 11.5, True, VIOLET)
estr = [
    ("TarifaDiaria", "precio_dia × dias", CYAN),
    ("TarifaSemanal", "bloques de 7 dias con descuento;\ndias sueltos a tarifa diaria",
     VIOLET),
    ("TarifaMensual", "bloques de 30 dias con descuento;\ndias sueltos a tarifa diaria",
     GOLD),
]
yy = 3.48
for nombre, f, c in estr:
    rect(s, 0.72, yy, 5.7, 0.78, WHITE, LINE, 0.75)
    rect(s, 0.72, yy, 0.055, 0.78, c)
    txt(s, 0.98, yy + 0.13, 2.4, 0.28, nombre, 13, True, NAVY, font=MONO)
    txt(s, 3.4, yy + 0.11, 2.9, 0.6, f, 11.5, False, MUTED, spacing=1.15)
    yy += 0.87
txt(s, 6.85, 1.6, 5.75, 0.3, "INTERFAZ Y CONTEXT", 11.5, True, VIOLET)
code_box(s, 6.85, 1.95, 5.75, 2.15,
         "class CalculoPlan(ABC):            # Strategy\n"
         "    @abstractmethod\n"
         "    def calcular_total(self, precio_dia, dias, plan):\n"
         "        ...\n\n"
         "# Context: la reserva no conoce la formula\n"
         "def calcular_total(self):\n"
         "    return self.plan.calculo().calcular_total(...)", 11)
txt(s, 6.85, 4.3, 5.75, 0.3, "LA REGLA DE NEGOCIO, EN UNA SOLA CLASE", 11.5, True, GOLD)
rect(s, 6.85, 4.65, 5.75, 1.05, LIGHT, LINE, 0.75)
txt(s, 7.08, 4.82, 5.3, 0.8,
    "El descuento aplica solo sobre bloques completos: 10 dias en plan semanal se cobran "
    "como 7 dias con descuento mas 3 dias a tarifa plena. Esa regla vive unicamente en "
    "TarifaSemanal.", 12.5, False, INK, spacing=1.15)
rect(s, 0.72, 6.15, 11.9, 0.55, NAVY)
txt(s, 1.0, 6.28, 11.4, 0.3,
    "Diferencia con State: la estrategia se elige desde afuera y no cambia sola. Una reserva "
    "pasa sola de Solicitada a Propuesta; nunca 'pasa' de TarifaDiaria a TarifaSemanal.",
    13, True, WHITE)
notes(s, "Adelantarse a la pregunta State vs Strategy: misma estructura, distinta intencion. "
         "El cabo suelto que abre el siguiente patron: si Reserva no nombra ninguna tarifa, "
         "quien decide cual se usa.")

# ==========================================================================
# 12 - Factory Method
# ==========================================================================
s = slide_base("Factory Method — planes de renta", "Patron creacional")
txt(s, 0.72, 1.6, 5.7, 0.3, "QUIEN DECIDE QUE TARIFA SE USA", 11.5, True, GOLD)
rect(s, 0.72, 1.95, 5.7, 1.05, LIGHT, LINE, 0.75)
rect(s, 0.72, 1.95, 0.055, 1.05, GOLD)
txt(s, 0.98, 2.12, 5.3, 0.8,
    "PlanRenta es el Creator. Lee su propio discriminador calculo_codigo (DIARIA, SEMANAL, "
    "MENSUAL) y fabrica la Strategy concreta. El Creator conoce a los productos; Reserva no.",
    12.5, False, INK, spacing=1.15)
code_box(s, 0.72, 3.15, 5.7, 2.5,
         "class PlanRenta(models.Model):        # Creator\n"
         "    calculo_codigo = models.CharField(...)\n\n"
         "    _PRODUCTOS = {\n"
         "        'DIARIA':  TarifaDiaria,\n"
         "        'SEMANAL': TarifaSemanal,\n"
         "        'MENSUAL': TarifaMensual,\n"
         "    }\n\n"
         "    def calculo(self):                # factory method\n"
         "        return self._PRODUCTOS[self.calculo_codigo]()", 10.5)
txt(s, 6.85, 1.6, 5.75, 0.3, "POR QUE ESTA VARIANTE", 11.5, True, GOLD)
rect(s, 6.85, 1.95, 5.75, 1.55, LIGHT, GOLD, 1.0)
txt(s, 7.08, 2.12, 5.3, 1.25,
    "Es lo que el propio GoF llama factory method parametrizado: un solo metodo que recibe "
    "un identificador, en vez de una subclase por producto. Se eligio porque el "
    "discriminador ya existia como columna en la base; crear una subclase de PlanRenta por "
    "cada plan habria duplicado la tabla sin ganar nada.", 12.5, False, INK, spacing=1.18)
txt(s, 6.85, 3.7, 5.75, 0.3, "PUNTO DE EXTENSION", 11.5, True, GREEN)
ext = [
    ("1", "Escribir la clase TarifaQuincenal"),
    ("2", "Agregar una entrada en el diccionario _PRODUCTOS"),
    ("3", "Insertar una fila en la tabla plan_renta"),
]
yy = 4.05
for n, t in ext:
    rect(s, 6.85, yy, 0.42, 0.42, GREEN)
    txt(s, 6.85, yy + 0.08, 0.42, 0.28, n, 13, True, WHITE, PP_ALIGN.CENTER, font=MONO)
    txt(s, 7.45, yy + 0.1, 5.1, 0.3, t, 13, False, INK)
    yy += 0.55
rect(s, 6.85, 5.75, 5.75, 0.62, NAVY)
txt(s, 7.08, 5.9, 5.3, 0.35,
    "Sin tocar Reserva ni las vistas. El mismo mecanismo se usa en EstadoReserva.crear().",
    12.5, True, WHITE)
rect(s, 0.72, 5.75, 5.7, 0.62, LIGHT, LINE, 0.75)
txt(s, 0.95, 5.9, 5.3, 0.35,
    "Una tabla no puede guardar un objeto que solo tiene comportamiento.", 12.5, True, NAVY)
notes(s, "Decir el matiz de la variante parametrizada antes de que lo pregunten: queda mejor "
         "que si lo sacan ellos. La pregunta 'como agregarian un plan nuevo' es la mejor que "
         "pueden hacer: es la demostracion de que el patron sirvio.")

# ==========================================================================
# 13 - Colaboracion
# ==========================================================================
s = slide_base("Los tres patrones trabajando juntos", "Que pasa cuando el cliente acepta")
txt(s, 0.72, 1.6, 11.9, 0.32,
    "Secuencia real de una sola operacion: reserva.aceptar(). Todo corre dentro de "
    "transaction.atomic con select_for_update.", 13.5, False, MUTED)
pasos = [
    ("1", "Cliente acepta", "La vista invoca reserva.aceptar()", NAVY),
    ("2", "State valida", "El estado Propuesta confirma que la transicion es legal", CYAN),
    ("3", "Factory fabrica", "plan.calculo() devuelve la Strategy segun el discriminador",
     GOLD),
    ("4", "Strategy calcula", "TarifaSemanal calcula el precio_total", VIOLET),
    ("5", "Se cierra el ciclo", "Reserva pasa a ACEPTADA y el vehiculo a RESERVADO", GREEN),
]
yy = 2.1
for n, t, d, c in pasos:
    rect(s, 0.72, yy, 8.6, 0.72, WHITE, LINE, 0.75)
    rect(s, 0.72, yy, 0.055, 0.72, c)
    rect(s, 1.0, yy + 0.15, 0.42, 0.42, c)
    txt(s, 1.0, yy + 0.23, 0.42, 0.28, n, 13, True, WHITE, PP_ALIGN.CENTER, font=MONO)
    txt(s, 1.62, yy + 0.11, 2.6, 0.3, t, 14, True, NAVY)
    txt(s, 4.3, yy + 0.13, 4.85, 0.42, d, 12.5, False, MUTED, spacing=1.1)
    yy += 0.8
rect(s, 9.6, 2.1, 3.02, 3.62, NAVY)
txt(s, 9.85, 2.32, 2.6, 0.3, "LA FRASE CLAVE", 11, True, CYAN)
rect(s, 9.85, 2.66, 0.7, 0.04, GOLD)
frases = [
    ("State", "decide SI se puede\ny QUE sigue", CYAN),
    ("Factory", "decide QUIEN calcula", GOLD),
    ("Strategy", "decide COMO se calcula", VIOLET),
]
yy = 2.9
for t, d, c in frases:
    txt(s, 9.85, yy, 2.55, 0.28, t, 15, True, c)
    txt(s, 9.85, yy + 0.28, 2.55, 0.5, d, 12, False, WHITE, spacing=1.15)
    yy += 0.83
rect(s, 9.85, 5.28, 2.55, 0.02, RGBColor(0x2E, 0x4A, 0x68))
txt(s, 9.6, 5.42, 3.02, 0.28, "Ninguno invade al otro.", 13, True, CYAN, PP_ALIGN.CENTER)
rect(s, 0.72, 6.2, 8.6, 0.5, LIGHT, LINE, 0.75)
txt(s, 0.95, 6.31, 8.2, 0.3,
    "Una transicion ilegal lanza ValidationError con mensaje concreto; el dato nunca se "
    "corrompe.", 12.5, True, NAVY)
notes(s, "Narrar la secuencia paso a paso. Rematar con la frase de la derecha, que es lo que "
         "se quiere que se lleven de la defensa.")

# ==========================================================================
# 14 - Arquitectura del demo
# ==========================================================================
s = slide_base("Arquitectura del demo", "Instruccion 2 — frontend, backend y persistencia")
capas = [
    ("Template", "Frontend", "Plantillas HTML, catalogo, formularios y dashboards", CYAN),
    ("View", "Backend", "Atiende la peticion, valida permisos y delega al modelo", VIOLET),
    ("Model", "Dominio", "Entidades, patrones de diseno y reglas de negocio", GOLD),
    ("PostgreSQL", "Persistencia", "11 tablas de negocio con constraints del DDL", GREEN),
]
x = 0.72
for t, sub, d, c in capas:
    rect(s, x, 1.7, 2.72, 1.85, WHITE, LINE, 0.9)
    rect(s, x, 1.7, 2.72, 0.06, c)
    txt(s, x + 0.2, 1.94, 2.35, 0.26, sub.upper(), 10.5, True, c)
    txt(s, x + 0.2, 2.22, 2.4, 0.36, t, 19, True, NAVY)
    txt(s, x + 0.2, 2.75, 2.4, 0.7, d, 11.5, False, MUTED, spacing=1.15)
    if x < 9.8:
        arrow(s, x + 2.79, 2.5, 0.28, 0.24, LINE)
    x += 3.07
txt(s, 0.72, 3.85, 6.0, 0.3, "LAS CUATRO APPS DEL PROYECTO", 11.5, True, CYAN)
data = [
    ["App", "Responsabilidad"],
    ["config", "Arranque del proyecto: configuracion general, rutas principales y conexion "
     "a la base de datos"],
    ["usuarios", "Login y personas del sistema: cliente, agente y administrador"],
    ["vehiculos", "Catalogo: tipos, ficha del vehiculo, fotos y favoritos"],
    ["reservas", "Proceso de renta: planes, solicitudes, propuestas, estados y calculo de "
     "tarifas"],
]
table(s, 0.72, 4.2, 7.4, data, [1.8, 5.6], row_h=0.5, head_h=0.38, size=12,
      first_bold=True)
txt(s, 8.45, 3.85, 4.17, 0.3, "DECISION CLAVE", 11.5, True, GOLD)
rect(s, 8.45, 4.2, 4.17, 2.5, LIGHT, LINE, 0.75)
rect(s, 8.45, 4.2, 0.055, 2.5, GOLD)
txt(s, 8.72, 4.4, 3.75, 2.1,
    "Los patrones viven en el modelo, no en la vista.\n\n"
    "La vista de aceptar una reserva tiene una sola linea de negocio: reserva.aceptar(). "
    "El resto son permisos y mensajes.\n\n"
    "Si manana cambia la regla de que estado puede aceptar, esa vista no se toca.",
    12.5, False, INK, spacing=1.2)
notes(s, "MTV = Model, Template, View. Insistir en que la vista no conoce ningun estado ni "
         "ninguna tarifa. Si piden ver codigo: reservas/views.py y reservas/estados.py.")

# ==========================================================================
# 15 - Modelo de datos
# ==========================================================================
s = slide_base("Modelo de datos", "11 tablas de negocio")
grupos = [
    ("Usuarios y roles", ["usuario", "administrador", "cliente", "agente"], CYAN,
     "Herencia por extension: usuario mas tabla de rol, en vez de una sola tabla ancha"),
    ("Catalogo", ["tipo_vehiculo", "vehiculo", "foto_vehiculo", "favorito"], VIOLET,
     "Ficha del vehiculo, estado operativo, imagenes y relacion cliente-vehiculo"),
    ("Renta", ["plan_renta", "reserva", "propuesta_renta"], GOLD,
     "plan_renta trae el discriminador del Factory; reserva es el Context del State"),
]
x = 0.72
for t, tablas, c, d in grupos:
    rect(s, x, 1.7, 3.9, 3.35, WHITE, LINE, 0.9)
    rect(s, x, 1.7, 3.9, 0.06, c)
    txt(s, x + 0.25, 1.95, 3.4, 0.32, t, 16, True, NAVY)
    yy = 2.42
    for tb in tablas:
        rect(s, x + 0.25, yy, 3.4, 0.36, LIGHT)
        txt(s, x + 0.42, yy + 0.07, 3.1, 0.26, tb, 12, True, CYAN_D, font=MONO)
        yy += 0.44
    txt(s, x + 0.25, 4.4, 3.4, 0.55, d, 11.5, False, MUTED, spacing=1.15)
    x += 4.06
txt(s, 0.72, 5.25, 6.0, 0.3, "DECISIONES DE DISENO", 11.5, True, CYAN)
bullets(s, 0.72, 5.6, 6.0, 1.1, [
    "Estados de vehiculo y reserva como VARCHAR con CHECK, alineados al patron State",
    "Snapshot de precios en reserva: aceptada una renta, no se recalcula con tarifas futuras",
], 12.5)
txt(s, 7.05, 5.25, 5.55, 0.3, "EL DDL MANDA", 11.5, True, GOLD)
rect(s, 7.05, 5.6, 5.55, 1.1, LIGHT, LINE, 0.75)
txt(s, 7.28, 5.75, 5.1, 0.85,
    "El esquema se diseno primero en SQL. Los modelos usan managed = False: las migraciones "
    "de negocio no reescriben el esquema, solo lo reflejan. Django agrego 7 tablas propias "
    "de infraestructura.", 12, False, INK, spacing=1.15)
notes(s, "El diagrama entidad-relacion completo esta en docs/MANUAL_TECNICO.md seccion 10 "
         "y en docs/diagramas/modelo_datos.md.")

# ==========================================================================
# 16 - Requerimientos funcionales
# ==========================================================================
s = slide_base("Requerimientos funcionales", "14 RF con plantilla IEEE 830 / ISO 29148")
txt(s, 0.72, 1.58, 11.9, 0.3,
    "Cada RF se documento con ID, funcion, descripcion, entradas, fuentes, salida, destino, "
    "accion, precondicion, postcondicion y prioridad, y se asocio a su caso de uso.",
    13, False, MUTED)
data = [
    ["ID", "Funcion", "Actor", "CDU", "Prioridad"],
    ["RF-001", "Registro de cliente", "Cliente", "01", "Alta"],
    ["RF-002", "Autenticacion de usuarios", "Todos", "02", "Alta"],
    ["RF-003", "Actualizacion de datos del cliente", "Cliente", "03", "Media"],
    ["RF-004", "Busqueda y consulta de vehiculos", "Cliente", "04", "Alta"],
    ["RF-005", "Gestion de favoritos", "Cliente", "05", "Media"],
    ["RF-006", "Solicitud de reserva", "Cliente", "06", "Alta"],
    ["RF-007", "Respuesta del cliente a la propuesta", "Cliente", "07", "Alta"],
]
table(s, 0.72, 2.05, 5.85, data, [1.0, 3.0, 1.1, 0.6, 1.0], row_h=0.36, head_h=0.38,
      size=11.5)
data2 = [
    ["ID", "Funcion", "Actor", "CDU", "Prioridad"],
    ["RF-008", "Gestion de vehiculos", "Agente", "08", "Alta"],
    ["RF-009", "Gestion de tipos de vehiculo", "Agente", "09", "Media"],
    ["RF-010", "Dashboard del agente", "Agente", "10", "Alta"],
    ["RF-011", "Propuesta de horario de entrega", "Agente", "11", "Alta"],
    ["RF-012", "Creacion de agentes", "Admin", "12", "Alta"],
    ["RF-013", "Baja de agentes", "Admin", "13", "Alta"],
    ["RF-014", "Calculo de tarifas segun plan", "Sistema", "07", "Alta"],
]
table(s, 6.75, 2.05, 5.87, data2, [1.0, 3.0, 1.1, 0.6, 1.0], row_h=0.36, head_h=0.38,
      size=11.5)
rect(s, 0.72, 5.0, 11.9, 1.65, LIGHT, LINE, 0.75)
rect(s, 0.72, 5.0, 0.055, 1.65, CYAN)
txt(s, 1.0, 5.18, 3.0, 0.3, "EJEMPLO — RF-007", 11.5, True, CYAN)
txt(s, 1.0, 5.48, 11.4, 1.05,
    "Accion:  aplicar la transicion de estado correspondiente; al aceptar, calcular la "
    "tarifa y marcar el vehiculo como RESERVADO; al rechazar, exigir motivo y dejar el "
    "vehiculo DISPONIBLE.\n"
    "Precondicion:  la reserva debe estar en estado PROPUESTA con horario propuesto.        "
    "Postcondicion:  la reserva queda en ACEPTADA o RECHAZADA.",
    12.5, False, INK, spacing=1.25)
notes(s, "No leer las tablas. Mostrar el nivel de detalle con el ejemplo de RF-007, que es el "
         "requerimiento donde se cruzan los tres patrones.")

# ==========================================================================
# 17 - Casos de uso de alto nivel
# ==========================================================================
s = slide_base("Casos de uso de alto nivel", "15 casos de uso, agrupados por modulo")
mods = [
    ("Modulo Cliente", CYAN, [
        ("CDU 01", "Registrar cliente"), ("CDU 02", "Iniciar sesion"),
        ("CDU 03", "Actualizar datos personales"), ("CDU 04", "Buscar y consultar vehiculos"),
        ("CDU 05", "Gestionar favoritos"), ("CDU 06", "Solicitar reserva"),
        ("CDU 07", "Aceptar o rechazar propuesta"),
    ]),
    ("Modulo Agente", VIOLET, [
        ("CDU 08", "Gestionar vehiculos"), ("CDU 09", "Gestionar tipos de vehiculo"),
        ("CDU 10", "Consultar dashboard"), ("CDU 11", "Proponer horario de entrega"),
        ("CDU 14", "Entregar vehiculo"), ("CDU 15", "Finalizar renta"),
    ]),
    ("Modulo Administrador", GOLD, [
        ("CDU 12", "Crear agente"), ("CDU 13", "Dar de baja agente"),
    ]),
]
x = 0.72
for t, c, casos in mods:
    rect(s, x, 1.68, 3.9, 4.2, WHITE, LINE, 0.9)
    rect(s, x, 1.68, 3.9, 0.06, c)
    txt(s, x + 0.25, 1.92, 3.4, 0.32, t, 15.5, True, NAVY)
    txt(s, x + 0.25, 2.26, 3.4, 0.26, f"{len(casos)} casos de uso", 11, True, c)
    yy = 2.62
    for cid, nombre in casos:
        txt(s, x + 0.25, yy, 0.85, 0.26, cid, 11, True, c, font=MONO)
        txt(s, x + 1.15, yy, 2.6, 0.26, nombre, 11.5, False, INK)
        yy += 0.42
    x += 4.06
rect(s, 0.72, 6.05, 11.9, 0.62, NAVY)
txt(s, 1.0, 6.2, 11.4, 0.32,
    "Los diagramas UML por modulo y los 15 casos de uso expandidos, con flujo normal y flujos "
    "alternativos, estan en el manual tecnico (secciones 7 y 8).", 13, True, WHITE)
notes(s, "El enunciado pide casos de uso de alto nivel y expandidos. Aqui se muestra el "
         "catalogo; la siguiente lamina ensena la profundidad de uno.")

# ==========================================================================
# 18 - CDU expandido
# ==========================================================================
s = slide_base("Caso de uso expandido", "CDU 07 — Aceptar o rechazar propuesta")
data = [
    ["Campo", "Contenido"],
    ["Actor principal", "Cliente"],
    ["Precondicion", "La reserva esta en estado PROPUESTA y tiene horario propuesto por el "
     "agente"],
    ["Postcondicion", "La reserva queda en ACEPTADA con precio calculado, o en RECHAZADA con "
     "motivo registrado"],
]
table(s, 0.72, 1.68, 11.9, data, [2.0, 9.9], row_h=0.44, head_h=0.36, size=12,
      first_bold=True)
txt(s, 0.72, 3.52, 5.7, 0.3, "FLUJO NORMAL", 11.5, True, GREEN)
flujo = [
    "El cliente abre la propuesta pendiente y revisa horario y plan",
    "Selecciona aceptar",
    "El sistema valida la transicion con el patron State",
    "Obtiene la Strategy del plan via Factory Method y calcula el total",
    "Guarda el snapshot de precios y marca el vehiculo como RESERVADO",
    "Muestra la confirmacion con el total calculado",
]
yy = 3.88
for i, p in enumerate(flujo, 1):
    txt(s, 0.72, yy, 0.35, 0.26, f"{i}.", 12, True, GREEN, font=MONO)
    txt(s, 1.12, yy, 5.3, 0.28, p, 12, False, INK)
    yy += 0.37
txt(s, 6.85, 3.52, 5.75, 0.3, "FLUJOS ALTERNATIVOS", 11.5, True, RED)
alt = [
    ("2a", "El cliente selecciona rechazar", "El sistema exige un motivo; la reserva pasa a "
     "RECHAZADA y el vehiculo vuelve a DISPONIBLE"),
    ("2b", "Motivo vacio en el rechazo", "Se bloquea la operacion; lo garantiza tambien el "
     "constraint chk_motivo_si_rechazada"),
    ("3a", "La reserva ya no esta en PROPUESTA", "El State lanza ValidationError con mensaje "
     "concreto y la vista lo muestra"),
]
yy = 3.88
for cid, t, d in alt:
    rect(s, 6.85, yy, 5.75, 0.83, LIGHT, LINE, 0.7)
    rect(s, 6.85, yy, 0.055, 0.83, RED)
    txt(s, 7.08, yy + 0.11, 0.5, 0.26, cid, 11.5, True, RED, font=MONO)
    txt(s, 7.6, yy + 0.1, 4.9, 0.28, t, 12, True, NAVY)
    txt(s, 7.6, yy + 0.4, 4.9, 0.4, d, 11, False, MUTED, spacing=1.1)
    yy += 0.9
rect(s, 0.72, 6.15, 5.7, 0.52, NAVY)
txt(s, 0.95, 6.27, 5.3, 0.3, "Todo el flujo corre dentro de una transaccion atomica.",
    12, True, WHITE)
notes(s, "Este CDU se eligio porque es donde colaboran los tres patrones y donde el RNF-04 "
         "de consistencia de estados se vuelve verificable.")

# ==========================================================================
# 19 - RNF
# ==========================================================================
s = slide_base("Requerimientos no funcionales", "10 RNF con metrica y criterio de aceptacion")
txt(s, 0.72, 1.58, 11.9, 0.3,
    "Cada RNF se describio de forma clara, medible y verificable, y se ancla a los casos de "
    "uso expandidos donde se comprueba.", 13, False, MUTED)
data = [
    ["ID", "Atributo de calidad", "Criterio de aceptacion medible", "Prioridad"],
    ["RNF-01", "Integridad y unicidad",
     "Correo y placa unicos; toda violacion de un CHECK produce error de validacion", "Alta"],
    ["RNF-02", "Seguridad de autenticacion",
     "Contrasena en hash; el login fallido no revela si el correo existe; CSRF activo",
     "Alta"],
    ["RNF-03", "Control de acceso por rol",
     "Un cliente no puede crear agentes ni modificar vehiculos de otro agente", "Alta"],
    ["RNF-04", "Consistencia de estados",
     "El 100 % de los intentos de transicion ilegal termina en error controlado", "Alta"],
    ["RNF-05", "Tiempo de respuesta",
     "Login, listado y creacion de reserva bajo 3 segundos con menos de 1 000 vehiculos",
     "Alta"],
    ["RNF-06", "Disponibilidad de datos",
     "Sin PostgreSQL la aplicacion falla de forma controlada y no corrompe el esquema",
     "Alta"],
    ["RNF-07", "Exactitud del calculo",
     "10 dias en plan semanal con 10 % de descuento coincide con la formula documentada",
     "Alta"],
    ["RNF-08", "Mantenibilidad",
     "Agregar un plan implica una Strategy y una entrada en el Factory, sin tocar el Context",
     "Media"],
    ["RNF-09", "Usabilidad basica",
     "Un usuario de prueba completa CDU 01, 06 y 07 con instrucciones minimas", "Media"],
    ["RNF-10", "Trazabilidad de propuestas",
     "Tras CDU 11 existe una fila historica en propuesta_renta ligada a la reserva", "Media"],
]
table(s, 0.72, 1.98, 11.9, data, [0.95, 2.6, 7.3, 1.05], row_h=0.4, head_h=0.36, size=10.5,
      first_bold=True, head_size=11)
txt(s, 0.72, 6.45, 11.9, 0.3,
    "Cada criterio se redacto para poder verificarse en el demo; la ficha completa de cada "
    "RNF esta en la seccion 9 del manual tecnico.", 11, False, MUTED)
notes(s, "El punto que evalua el auxiliar: que cada RNF tenga metrica, no adjetivos. Senalar "
         "RNF-04 y RNF-07, que son los que verifican directamente los patrones.")

# ==========================================================================
# 20 - Git Flow
# ==========================================================================
s = slide_base("Estrategia de branching", "Instruccion 3 — Git Flow")
rect(s, 0.72, 1.68, 11.9, 1.5, LIGHT, LINE, 0.75)
chip(s, 1.0, 2.0, 1.9, 0.5, "feature/*", VIOLET, WHITE, 13)
arrow(s, 3.05, 2.11, 0.85, 0.28, LINE)
txt(s, 3.0, 1.76, 1.0, 0.24, "merge", 10.5, False, MUTED, PP_ALIGN.CENTER)
chip(s, 4.05, 2.0, 1.9, 0.5, "develop", CYAN, WHITE, 13)
arrow(s, 6.1, 2.11, 0.85, 0.28, LINE)
txt(s, 6.05, 1.76, 1.0, 0.24, "release", 10.5, False, MUTED, PP_ALIGN.CENTER)
chip(s, 7.1, 2.0, 1.9, 0.5, "main", GREEN, WHITE, 13)
chip(s, 9.4, 2.0, 1.9, 0.5, "hotfix/*", RED, WHITE, 13)
txt(s, 9.4, 2.58, 3.0, 0.24, "nace de main, vuelve a main y develop", 10.5, False, MUTED)
txt(s, 1.0, 2.62, 7.9, 0.3,
    "Cada feature nace de develop y regresa a develop. Cuando varias features quedan "
    "consolidadas, develop se fusiona a main.", 12, False, INK)
txt(s, 0.72, 3.35, 5.7, 0.3, "CONVENCIONES DEL EQUIPO", 11.5, True, CYAN)
data = [
    ["Elemento", "Convencion"],
    ["Rama feature", "feature/<nombre>-<carne>"],
    ["Rama hotfix", "hotfix/<descripcion>-<carne>"],
    ["Commit", "#carne: mensaje"],
    ["Ejemplo real", "201900131: feat(vehiculos) catalogo y flota"],
]
table(s, 0.72, 3.7, 5.7, data, [1.9, 3.8], row_h=0.42, head_h=0.38, size=12,
      first_bold=True)
txt(s, 6.85, 3.35, 5.75, 0.3, "EVIDENCIA EN EL REPOSITORIO", 11.5, True, GREEN)
ramas = [
    "feature/setup-django_201905884",
    "feature/database-202112395",
    "feature/usuarios-models-201905884",
    "feature/vehiculos-models-201905884",
    "feature/reservas-models-201905884",
    "feature/design-patterns-201905884",
    "feature/cliente-agente-201801521",
    "feature/frontend-base-201900131",
    "feature/documentacionv3-201701078",
]
yy = 3.7
for r in ramas:
    txt(s, 6.85, yy, 5.75, 0.26, "•  " + r, 11.5, False, CYAN_D, font=MONO)
    yy += 0.28
rect(s, 0.72, 5.85, 5.7, 0.85, NAVY)
txt(s, 0.95, 6.0, 5.3, 0.6,
    "Las ramas de feature no se eliminaron, tal como exige el enunciado. Los cinco "
    "integrantes tienen commits validos.", 12.5, True, WHITE, spacing=1.15)
notes(s, "Mostrar en vivo git branch -a y git log --oneline si lo piden. Cada rama lleva el "
         "carne de quien la trabajo, lo que hace trazable la participacion individual.")

# ==========================================================================
# 21 - Entregables
# ==========================================================================
s = slide_base("Entregables y documentacion", "Lo que se sube a UEDI")
ent = [
    ("Repositorio privado", "AYD2_A_2S2026_PRACTICA1_G2",
     "Con el auxiliar de la seccion A (kjhg4589) agregado como colaborador", CYAN),
    ("Manual tecnico", "docs/MANUAL_TECNICO.md",
     "Formato markdown dentro del repositorio, como exige el enunciado", VIOLET),
    ("Diagramas UML", "docs/diagramas/",
     "Patrones (clases, estados y secuencia) y modelo entidad-relacion", GOLD),
    ("Demo funcional", "Django + PostgreSQL",
     "Frontend, backend y persistencia con datos de demostracion cargados", GREEN),
]
x = 0.72
for t, ref, d, c in ent:
    rect(s, x, 1.7, 2.9, 2.35, WHITE, LINE, 0.9)
    rect(s, x, 1.7, 2.9, 0.06, c)
    txt(s, x + 0.22, 1.98, 2.5, 0.32, t, 14.5, True, NAVY)
    rect(s, x + 0.22, 2.42, 2.5, 0.34, LIGHT)
    txt(s, x + 0.33, 2.5, 2.35, 0.28, ref, 10.5, True, CYAN_D, font=MONO)
    txt(s, x + 0.22, 2.95, 2.5, 0.95, d, 11.5, False, MUTED, spacing=1.18)
    x += 3.06
txt(s, 0.72, 4.3, 6.0, 0.3, "CONTENIDO DEL MANUAL TECNICO", 11.5, True, CYAN)
secciones = [
    "1. Introduccion", "2. Descripcion del sistema", "3. Metodologia Git Flow",
    "4. Seleccion de framework", "5. Patrones de diseno", "6. Requerimientos funcionales",
    "7. Casos de uso de alto nivel", "8. Casos de uso expandidos",
    "9. Requerimientos no funcionales", "10. Modelo de datos",
    "11. Arquitectura de la solucion", "12. Referencias",
]
x, yy = 0.72, 4.68
for i, sec in enumerate(secciones):
    rect(s, x, yy, 2.9, 0.36, LIGHT, LINE, 0.6)
    txt(s, x + 0.15, yy + 0.06, 2.7, 0.26, sec, 11.5, False, INK)
    x += 3.06
    if (i + 1) % 4 == 0:
        x, yy = 0.72, yy + 0.44
rect(s, 0.72, 6.15, 11.9, 0.55, NAVY)
txt(s, 1.0, 6.28, 11.4, 0.3,
    "Participacion verificable: cada integrante tiene al menos un commit valido, con "
    "contenido real, bajo el formato #carne: mensaje.", 13, True, WHITE)
notes(s, "Recordar que el enunciado exige markdown y solo entrega por UEDI. Confirmar que el "
         "auxiliar ya fue agregado al repositorio privado.")

# ==========================================================================
# 22 - Demo
# ==========================================================================
s = slide_base("Recorrido de la demostracion", "Que se va a mostrar en vivo")
demo = [
    ("1", "Registro y login", "Cliente", "Alta con correo unico, licencia y fotografia; "
     "ingreso con credenciales", CYAN),
    ("2", "Catalogo y favoritos", "Cliente", "Busqueda de vehiculos, ficha completa y "
     "marcado de favoritos", CYAN),
    ("3", "Solicitud de reserva", "Cliente", "Fecha, hora, direccion de entrega y cantidad "
     "de dias → estado SOLICITADA", CYAN),
    ("4", "Dashboard y propuesta", "Agente", "Vehiculos disponibles, reservados y en renta; "
     "propuesta de horario → PROPUESTA", VIOLET),
    ("5", "Aceptacion y calculo", "Cliente", "Aceptar dispara State + Factory + Strategy y "
     "fija el precio total", GOLD),
    ("6", "Entrega y cierre", "Agente", "EN_RENTA y luego FINALIZADA, con el vehiculo "
     "sincronizado", GREEN),
    ("7", "Gestion de agentes", "Admin", "Alta y baja de agentes de renta", RED),
]
yy = 1.72
for n, t, actor, d, c in demo:
    rect(s, 0.72, yy, 11.9, 0.62, WHITE, LINE, 0.7)
    rect(s, 0.72, yy, 0.055, 0.62, c)
    rect(s, 1.0, yy + 0.12, 0.38, 0.38, c)
    txt(s, 1.0, yy + 0.18, 0.38, 0.26, n, 12, True, WHITE, PP_ALIGN.CENTER, font=MONO)
    txt(s, 1.58, yy + 0.17, 2.7, 0.3, t, 13.5, True, NAVY)
    chip(s, 4.4, yy + 0.14, 1.0, 0.34, actor, c, WHITE, 10)
    txt(s, 5.65, yy + 0.18, 6.8, 0.3, d, 12, False, MUTED)
    yy += 0.7
notes(s, "Probar el recorrido completo antes de exponer. Si preguntan por el codigo, abrir "
         "reservas/estados.py: es el archivo que mejor se defiende solo.")

# ==========================================================================
# 23 - Preguntas anticipadas
# ==========================================================================
s = slide_base("Preguntas anticipadas", "Lo que conviene tener resuelto")
qa = [
    ("Por que tres patrones si se pedia uno?",
     "El dominio los pidio: ciclo de vida y calculo de precio son problemas distintos. Con "
     "solo State se cumplia, pero el precio habria quedado con condicionales dentro de "
     "Reserva.", CYAN),
    ("Ese Factory Method es realmente el del GoF?",
     "Es la variante parametrizada que el propio GoF describe: un metodo que recibe un "
     "identificador. Se eligio porque el discriminador ya existia como columna en la base.",
     GOLD),
    ("Que pasa con una transicion ilegal?",
     "Se lanza ValidationError con mensaje concreto, por ejemplo 'no se puede aceptar una "
     "reserva en estado EN_RENTA'. Todo corre en transaction.atomic, el dato no se corrompe.",
     RED),
    ("Como agregarian un plan nuevo?",
     "Tres pasos sin tocar Reserva ni las vistas: una clase TarifaQuincenal, una entrada en "
     "_PRODUCTOS y una fila en plan_renta.", GREEN),
    ("State y Strategy no son lo mismo?",
     "Misma estructura, distinta intencion: en State el objeto esta en una situacion y el "
     "patron decide la siguiente; en Strategy usa un algoritmo elegido desde afuera.",
     VIOLET),
    ("Por que no guardaron el objeto de estado en la base?",
     "No tiene datos propios, solo comportamiento. Se guarda el codigo y se reconstruye el "
     "objeto con EstadoReserva.crear().", NAVY),
]
x, yy = 0.72, 1.68
for i, (q, a, c) in enumerate(qa):
    rect(s, x, yy, 5.85, 1.55, WHITE, LINE, 0.8)
    rect(s, x, yy, 0.055, 1.55, c)
    txt(s, x + 0.25, yy + 0.16, 5.4, 0.3, q, 13.5, True, NAVY)
    txt(s, x + 0.25, yy + 0.55, 5.4, 0.9, a, 11.5, False, MUTED, spacing=1.18)
    if i % 2 == 0:
        x += 6.05
    else:
        x, yy = 0.72, yy + 1.68
notes(s, "Lamina de red de seguridad. No leerla: pasarla rapido y volver si preguntan.")

# ==========================================================================
# 24 - Cierre
# ==========================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 0, 0, 13.333, 0.22, CYAN)
rect(s, 0, 7.28, 13.333, 0.22, CYAN)
txt(s, 0.9, 1.05, 11.5, 0.32, "CONCLUSIONES", 12.5, True, CYAN)
rect(s, 0.9, 1.45, 1.05, 0.06, GOLD)
txt(s, 0.9, 1.75, 11.5, 0.6, "Lo que deja la practica", 34, True, WHITE)
cierres = [
    ("Framework", "Django 6.0.7 con PostgreSQL cubrio frontend, backend, ORM y "
     "autenticacion en un solo stack, dejando espacio para los patrones en el dominio."),
    ("Patrones", "State, Strategy y Factory Method colaboran en una sola operacion: cero "
     "condicionales sobre el estado dentro de Reserva."),
    ("Extensibilidad", "Agregar un plan o un estado es escribir una clase y registrarla; "
     "las clases existentes no se modifican."),
    ("Proceso", "Git Flow con main, develop y ramas feature por integrante, con commits "
     "trazables bajo el formato #carne: mensaje."),
]
yy = 2.68
for t, d in cierres:
    rect(s, 0.9, yy, 0.05, 0.8, CYAN)
    txt(s, 1.2, yy + 0.02, 2.4, 0.32, t, 16, True, CYAN)
    txt(s, 3.7, yy + 0.03, 8.7, 0.7, d, 13.5, False, RGBColor(0xD4, 0xE2, 0xEE),
        spacing=1.18)
    yy += 0.95
rect(s, 0.9, 6.5, 11.5, 0.02, RGBColor(0x2E, 0x4A, 0x68))
txt(s, 0.9, 6.7, 7.0, 0.3, "Gracias. Preguntas?", 16, True, WHITE)
txt(s, 8.5, 6.72, 3.9, 0.3, "Grupo 2  ·  Seccion A  ·  2S 2026", 12, False,
    RGBColor(0x9F, 0xB4, 0xC8), PP_ALIGN.RIGHT)
notes(s, "Cerrar con el beneficio concreto, no con teoria. Abrir a preguntas con la lamina "
         "anterior a mano.")


# --------------------------------------------------------------------------
out = Path(__file__).with_name("Presentacion_Practica1_AutoRent_G2.pptx")
prs.save(out)
print(f"OK  {len(prs.slides.__iter__.__self__._sldIdLst)} diapositivas -> {out}")
